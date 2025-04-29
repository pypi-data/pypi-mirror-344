import csv
import os
import pandas as pd
import re
import docx
import fitz
import uuid
from typing import Dict,Optional
from pptx import Presentation
from huggingface_hub import InferenceClient
from dotenv import load_dotenv
from openai import OpenAI
from global_parser.file_parser.type.image import process_image_with_pixtral
from global_parser.utils.upload import upload_bytes_to_s3
from bs4 import BeautifulSoup
import pdfplumber
import requests
import tempfile
load_dotenv()
HF_TOKEN = os.getenv("HF_TOKEN","")
VAIA_IMAGE_TO_TEXT_URL = os.getenv("VAIA_IMAGE_TO_TEXT_URL","")

client = InferenceClient(api_key=HF_TOKEN)
client_llama = OpenAI(
	base_url=VAIA_IMAGE_TO_TEXT_URL, 
	api_key=HF_TOKEN 
)
s3_bucket = "ft-vaia"
pixtral_api_key = os.getenv("PIXTAL_API_KEY","")

def parse_txt(file_path):
    with open(file_path, 'r') as file:
        text = file.read()
    return text
def parse_pdf(temp_file_path):
    

    text = ""
    print(temp_file_path,'temp_file_path')
    with pdfplumber.open(temp_file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text

def parse_pdf_with_links(file_path):
    text = ""
    links = set()  # Using set to avoid duplicate links
    
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            # Extract text from the page
            text += page.extract_text() + "\n"
            
            # Extract annotations (which include links)
            annots = page.hyperlinks
            if annots:
                for annot in annots:
                    if annot and 'uri' in annot:
                        link_url = annot['uri']
                        links.add(link_url)
    
    # Process each unique link
    if links:
        text += "\n=== Links Found in PDF ===\n"
        for link in links:
            text += f"Link: {link}\n"
            # Try to fetch and extract text from the link if it's accessible
            try:
                response = requests.get(link, timeout=10, verify=False)
                if response.status_code == 200:
                    # Determine file type from URL and content
                    is_pdf = link.lower().endswith('.pdf')
                    is_docx = link.lower().endswith('.docx')
                    is_xlsx = link.lower().endswith('.xlsx') or link.lower().endswith('.xls')
                    is_pptx = link.lower().endswith('.pptx')
                    
                    # Create a temporary file with appropriate extension
                    file_ext = '.pdf' if is_pdf else '.docx' if is_docx else '.xlsx' if is_xlsx else '.pptx' if is_pptx else '.html'
                    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
                        temp_file.write(response.content)
                        temp_file_path = temp_file.name
                    
                    try:
                        if is_pdf:
                            # Parse the linked PDF
                            with pdfplumber.open(temp_file_path) as linked_pdf:
                                linked_text = ""
                                for page in linked_pdf.pages:
                                    linked_text += page.extract_text() + "\n"
                                text += f"\nContent from PDF {link}:\n"
                                text += f"{linked_text}\n"
                        elif is_docx:
                            # Parse DOCX file
                            linked_text = parse_docx(temp_file_path)
                            text += f"\nContent from DOCX {link}:\n"
                            text += f"{linked_text}\n"
                        elif is_xlsx:
                            # Parse XLSX file - Note: We'll only extract text without template
                            df = pd.read_excel(temp_file_path)
                            linked_text = df.to_string()
                            text += f"\nContent from Excel {link}:\n"
                            text += f"{linked_text}\n"
                        elif is_pptx:
                            # Parse PPTX file
                            linked_text = parse_pptx(temp_file_path)
                            text += f"\nContent from PowerPoint {link}:\n"
                            text += f"{linked_text}\n"
                        else:
                            # Handle web content
                            soup = BeautifulSoup(response.text, 'html.parser')
                            link_text = soup.get_text(separator='\n', strip=True)
                            text += f"\nContent from {link}:\n"
                            text += f"{link_text}\n"
                        
                        text += "=" * 80 + "\n"
                    finally:
                        # Clean up temporary file
                        try:
                            os.unlink(temp_file_path)
                        except Exception as e:
                            print(f"Error deleting temporary file: {e}")
            except Exception as e:
                text += f"Could not process link content: {str(e)}\n"
                continue
    
    return text
def parse_docx(temp_file_path):
    doc = docx.Document(temp_file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    print("Text : " , text)
    return text

def parse_xlsx(file_path, input_template):
    print("Parsing Excel file across all sheets")
    
    # Read all sheets into a dictionary of DataFrames
    sheets = pd.read_excel(file_path, sheet_name=None)
    
    # Extract columns from the template
    columns_in_template = re.findall(r"\{(.*?)\}", input_template)
    print("Columns in template:", columns_in_template)
    
    # Initialize a list to store text summaries from all sheets
    all_text_summaries = []
    
    for sheet_name, df in sheets.items():
        print(f"Processing sheet: {sheet_name}")
        df = df.head(10)
        # Check for missing columns in the Excel sheet
        missing_columns = [col for col in columns_in_template if col not in df.columns]
        if missing_columns:
            print(f"Error in sheet '{sheet_name}': Missing columns: {', '.join(missing_columns)}")
            continue  # Skip this sheet
        
        df = df.astype(str)
        # Generate the "text" column using the template
        df['text'] = df.apply(lambda row: input_template.format(**row), axis=1)
        
        print(f"Sample text from sheet '{sheet_name}':")
        print(df['text'].head())  # Print the first few rows for debugging
        
        # Concatenate all rows from this sheet and append to the summary list
        all_text_summaries.append("\n".join(df['text'].tolist()))
    
    # Combine summaries from all sheets
    combined_summary = "\n\n---\n\n".join(all_text_summaries)
    
    print("Generated combined text summary:")
    print(combined_summary[:500])  # Print the first 500 characters for debugging
    
    return combined_summary

def parse_xlsx_without_template(file_path):
    df = pd.read_excel(file_path, engine='openpyxl')
    df = df.head(10)
    text = '\n'.join(df.apply(lambda row: ' '.join(row.values.astype(str)), axis=1))
    return text

def parse_csv(file_path, input):
    print("Parsing CSV")
    df = pd.read_csv(file_path)
    template = input
    # Extract columns from the template
    columns_in_template = re.findall(r"\{(.*?)\}", template)
    print("Columns in template: ", columns_in_template)

    # Check for missing columns
    missing_columns = [col for col in columns_in_template if col not in df.columns]
    if missing_columns:
        print(f"Error: Missing columns in CSV file: {', '.join(missing_columns)}")
        return f"Error: Missing columns in CSV file: {', '.join(missing_columns)}"
    else:
        # Generate the "text" column using the template
        df['text'] = df.apply(lambda row: template.format(**row), axis=1)

        # Combine all text into a single string
        combined_text = "\n".join(df['text'].tolist())

    print(f"Generated text from CSV: {combined_text}")
    return combined_text

def parse_csv_without_template(file_path):
    df = pd.read_csv(file_path)
    df = df.head(10)
    # text = df.to_string(index=False)
    header = ' '.join(df.columns.astype(str))
    rows = '\n'.join(df.apply(lambda row: ' '.join(row.values.astype(str)), axis=1))
    text = header + '\n' + rows
    print(text)
    return text
def text_to_csv(text: str, csv_filename: str = "output.csv") -> str:
    # lines = text.splitlines()
    sentences = text.replace("\n", " ").split(".")
    
    # Remove any empty strings resulting from the split
    sentences = [sentence.strip() for sentence in sentences if len(sentence.strip()) >= 10]
    
    csv_file_path = os.path.join("/tmp", csv_filename)

    with open(csv_file_path, mode="w", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["text"])  # Column header
        for line in sentences:
            writer.writerow([line])

    return csv_file_path

# def rich_text_to_csv(result: dict, csv_filename: str = "output.csv") -> str:
#     text = result["text_content"]
#     images_content = result["images_content"]

#     # Split the text_content into sentences, delimited by "."
#     sentences = text.replace("\n", " ").split(".")
#     sentences = [sentence.strip() for sentence in sentences if len(sentence.strip()) >= 10]  # Remove short and empty sentences
    
#     csv_file_path = os.path.join("/tmp", csv_filename)

#     with open(csv_file_path, mode="w", newline="") as file:
#         writer = csv.writer(file)
#         writer.writerow(["text"])  # Column header
        
#         # Write each sentence from text_content
#         for sentence in sentences:
#             writer.writerow([sentence])
        
#         # Leave an empty row to separate text_content and image descriptions
#         writer.writerow([])

#         # Write each image description below text_content
#         for image in images_content:
#             image_row = image["description"]
#             writer.writerow([image_row])

#     print(f"CSV file with formatted text saved as {csv_file_path}.")
#     return csv_file_path

def rich_text_to_csv(result: Dict, csv_filename: str = "output.csv") -> str:        
    # Check for required keys
    print("result inside rich_text_to_csv : " , result)
    # if "text_content" not in result or "images_content" not in result:
    #     raise KeyError("result dict must contain 'text_content' and 'images_content' keys")
        
    text = result["text_content"]
    images_content = result["images_content"]
    
    if not isinstance(text, str):
        raise TypeError("text_content must be a string")
    
    if not isinstance(images_content, list):
        raise TypeError("images_content must be a list")

    # Split the text_content into sentences, delimited by "."
    sentences = text.replace("\n", " ").split(".")
    sentences = [sentence.strip() for sentence in sentences if len(sentence.strip()) >= 10]

    os.makedirs("/tmp", exist_ok=True)
    
    csv_file_path = os.path.join("/tmp", csv_filename)

    try:
        with open(csv_file_path, mode="w", newline="", encoding="utf-8") as file:
            writer = csv.writer(file)
            writer.writerow(["text"]) 
            
            for sentence in sentences:
                writer.writerow([sentence])
            
            writer.writerow([])

            for image in images_content:
                if not isinstance(image, dict) or "description" not in image:
                    raise TypeError("Each image in images_content must be a dict with a 'description' key")
                image_row = image["description"]
                writer.writerow([image_row])

        print(f"CSV file with formatted text saved as {csv_file_path}.")
        return csv_file_path
        
    except OSError as e:
        raise OSError(f"Failed to write to CSV file: {e}")
    
def parse_pdf_with_images(temp_file_path):
    
    result = {
        "text_content": "",
        "images_content": [],
        "combined_content": ""
    }
    with pdfplumber.open(temp_file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                result["text_content"] += text + "\n"
    
    doc = fitz.open(temp_file_path)
    for page_num in range(doc.page_count):
        page = doc[page_num]
        image_list = page.get_images()

        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            # Generate unique filename for S3
            file_name = f"pdf-processed/pdf_image_{uuid.uuid4()}.{base_image['ext']}"
            
            # Upload to S3
            image_url = upload_bytes_to_s3(file_bytes=image_bytes,
                bucket=s3_bucket,
                object_name=file_name
            )
            
            if image_url:
                # Process image with Pixtral
                image_description = process_image_with_pixtral(image_url, pixtral_api_key)
                
                if image_description:
                    image_content = {
                        "page_number": page_num + 1,
                        "image_url": image_url,
                        "description": image_description
                    }
                    result["images_content"].append(image_content)
                    
                    # Add image description to the combined content
                    result["combined_content"] += f"\n[Image on page {page_num + 1}: {image_description}]\n"
    # print(result)
    # Combine all content
    result["combined_content"] = result["text_content"] + "\n" + \
        "\n".join([f"Image {i+1}: {img['description']}. " 
                  for i, img in enumerate(result["images_content"])])
    
    doc.close()
    # print(result)
    # return result
    text = result["combined_content"]
    return text

def parse_docx_with_images(temp_file_path):
    result = {
        "text_content": "",
        "images_content": [],
        "combined_content": ""
    }
    
    doc = docx.Document(temp_file_path)
    
    for para in doc.paragraphs:
        result["text_content"] += para.text + "\n"
    
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel.target_part.blob
            
            file_name = f"RAG-images-processed/docx_image_{uuid.uuid4()}.jpeg"
            
            image_url = upload_bytes_to_s3(
                file_bytes=image_data,
                bucket=s3_bucket,
                object_name=file_name
            )
            
            if image_url:
                image_description = process_image_with_pixtral(image_url, pixtral_api_key)
                
                if image_description:
                    page_number = find_image_paragraph_number(doc, rel.target_ref)
                    
                    image_content = {
                        "page_number": page_number,
                        "image_url": image_url,
                        "description": image_description
                    }
                    result["images_content"].append(image_content)
                    
                    # Add image description to the combined content
                    result["combined_content"] += f"\n[Image on paragraph {page_number}: {image_description}]\n"
    
    # Combine all content
    result["combined_content"] = result["text_content"] + "\n" + \
        "\n".join([f"[Image {i+1}: {img['description']}]" 
                  for i, img in enumerate(result["images_content"])])
    print("Result : " , result)   
    # return result
    return result["combined_content"]


def find_image_paragraph_number(doc, target_ref):
    for i, paragraph in enumerate(doc.paragraphs):
        for run in paragraph.runs:
            for content in run._element:
                if hasattr(content, 'graphic') and target_ref in str(content.graphic):
                    return i + 1
    return 0

def parse_pptx(temp_file_path):
    print("Parsing PPTX")
    prs = Presentation(temp_file_path)
    print("Parsing PPTX",prs,temp_file_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text


def download_and_extract_headers_xlsx(file_path: str) -> Dict[str, int]:
    # Read headers from the Excel file
    headers = {}
    df = pd.read_excel(file_path, sheet_name=0, engine='openpyxl')  # Read the first sheet
    headers = {header: idx for idx, header in enumerate(df.columns)}

    print("Headers with column numbers:", headers)
    # Join header keys with "and" word
    header_keys = list(headers.keys())
    header_keys = [f"{{{key}}}" for key in header_keys]
    joined_headers = " and ".join(header_keys)
    print("Joined headers:", joined_headers)
    return joined_headers

def parse_pdf_with_links_and_images(temp_file_path):
    """
    Parse a PDF file to extract text, links, and images with descriptions.
    
    Args:
        temp_file_path: Path to the PDF file
        
    Returns:
        Extracted text content including links and image descriptions
    """
    
    # Initialize result dictionary similar to parse_pdf_with_images
    result = {
        "text_content": "",
        "images_content": [],
        "combined_content": ""
    }
    
    # Extract text and links (from parse_pdf_with_links)
    text = ""
    links = set()  # Using set to avoid duplicate links
    
    with pdfplumber.open(temp_file_path) as pdf:
        for page in pdf.pages:
            # Extract text from the page
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
                result["text_content"] += page_text + "\n"
            
            # Extract annotations (which include links)
            annots = page.hyperlinks
            if annots:
                for annot in annots:
                    if annot and 'uri' in annot:
                        link_url = annot['uri']
                        links.add(link_url)
    
    # Process each unique link (from parse_pdf_with_links)
    if links:
        text += "\n=== Links Found in PDF ===\n"
        for link in links:
            text += f"Link: {link}\n"
            # Try to fetch and extract text from the link if it's accessible
            try:
                response = requests.get(link, timeout=10, verify=False)
                if response.status_code == 200:
                    # Determine file type from URL and content
                    is_pdf = link.lower().endswith('.pdf')
                    is_docx = link.lower().endswith('.docx')
                    is_xlsx = link.lower().endswith('.xlsx') or link.lower().endswith('.xls')
                    is_pptx = link.lower().endswith('.pptx')
                    
                    # Create a temporary file with appropriate extension
                    file_ext = '.pdf' if is_pdf else '.docx' if is_docx else '.xlsx' if is_xlsx else '.pptx' if is_pptx else '.html'
                    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
                        temp_file.write(response.content)
                        temp_link_file_path = temp_file.name
                    
                    try:
                        if is_pdf:
                            # Parse the linked PDF
                            with pdfplumber.open(temp_link_file_path) as linked_pdf:
                                linked_text = ""
                                for page in linked_pdf.pages:
                                    linked_text += page.extract_text() + "\n"
                                text += f"\nContent from PDF {link}:\n"
                                text += f"{linked_text}\n"
                        elif is_docx:
                            # Parse DOCX file
                            linked_text = parse_docx(temp_link_file_path)
                            text += f"\nContent from DOCX {link}:\n"
                            text += f"{linked_text}\n"
                        elif is_xlsx:
                            # Parse XLSX file - Note: We'll only extract text without template
                            df = pd.read_excel(temp_link_file_path)
                            linked_text = df.to_string()
                            text += f"\nContent from Excel {link}:\n"
                            text += f"{linked_text}\n"
                        elif is_pptx:
                            # Parse PPTX file
                            linked_text = parse_pptx(temp_link_file_path)
                            text += f"\nContent from PowerPoint {link}:\n"
                            text += f"{linked_text}\n"
                        else:
                            # Handle web content
                            soup = BeautifulSoup(response.text, 'html.parser')
                            link_text = soup.get_text(separator='\n', strip=True)
                            text += f"\nContent from {link}:\n"
                            text += f"{link_text}\n"
                        
                        text += "=" * 80 + "\n"
                    finally:
                        # Clean up temporary file
                        try:
                            os.unlink(temp_link_file_path)
                        except Exception as e:
                            print(f"Error deleting temporary file: {e}")
            except Exception as e:
                text += f"Could not process link content: {str(e)}\n"
                continue
    
    # Extract and process images (from parse_pdf_with_images)
    doc = fitz.open(temp_file_path)
    for page_num in range(doc.page_count):
        page = doc[page_num]
        image_list = page.get_images()

        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            # Generate unique filename for S3
            file_name = f"pdf-processed/pdf_image_{uuid.uuid4()}.{base_image['ext']}"
            
            # Upload to S3
            image_url = upload_bytes_to_s3(file_bytes=image_bytes,
                bucket=s3_bucket,
                object_name=file_name
            )
            
            if image_url:
                # Process image with Pixtral
                image_description = process_image_with_pixtral(image_url, pixtral_api_key)
                
                if image_description:
                    image_content = {
                        "page_number": page_num + 1,
                        "image_url": image_url,
                        "description": image_description
                    }
                    result["images_content"].append(image_content)
                    
                    # Add image description to text and combined content
                    text += f"\n[Image on page {page_num + 1}: {image_description}]\n"
                    result["combined_content"] += f"\n[Image on page {page_num + 1}: {image_description}]\n"
    
    doc.close()
    
    # Combine all content (from parse_pdf_with_images)
    result["combined_content"] = result["text_content"] + "\n" + \
        "\n".join([f"Image {i+1}: {img['description']}. " 
                  for i, img in enumerate(result["images_content"])])
    
    # Return text as in original function
    return text