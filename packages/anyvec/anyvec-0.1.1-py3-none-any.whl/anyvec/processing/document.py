import pdfplumber
import io
from docx import Document
import subprocess
import fitz
import base64


def extract_text_plain(buffer: bytes) -> str:
    return buffer.decode("utf-8")


def extract_text_pdf(buffer: bytes) -> str:
    with pdfplumber.open(io.BytesIO(buffer)) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
        return text


def extract_images_pdf(buffer: bytes) -> list[str]:
    images: list[str] = []
    pdf_document = fitz.open("pdf", buffer)
    for page in pdf_document:
        for _, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_buffer = base_image["image"]
            encoded_image_buffer = base64.b64encode(image_buffer).decode("utf-8")
            images.append(encoded_image_buffer)
    return images

def extract_images_docx(buffer: bytes) -> list[str]:
    from docx import Document
    import io
    images = []
    doc = Document(io.BytesIO(buffer))
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img_bytes = rel.target_part.blob
            images.append(base64.b64encode(img_bytes).decode("utf-8"))
    return images

def extract_images_pptx_like(buffer: bytes) -> list[str]:
    from pptx import Presentation
    import io
    images = []
    prs = Presentation(io.BytesIO(buffer))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                images.append(base64.b64encode(image.blob).decode("utf-8"))
    return images

def extract_images_odt(buffer: bytes) -> list[str]:
    from odf.opendocument import load
    from odf.draw import Frame, Image
    import tempfile
    images = []
    with tempfile.NamedTemporaryFile(delete=True, suffix=".odt") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        for frame in doc.getElementsByType(Frame):
            for img in frame.getElementsByType(Image):
                href = img.getAttribute("href")
                if href and href.startswith("Pictures/"):
                    for f in doc.Pictures:
                        if f == href:
                            img_bytes = doc.Pictures[f]
                            images.append(base64.b64encode(img_bytes).decode("utf-8"))
    return images

def extract_images_odp(buffer: bytes) -> list[str]:
    from odf.opendocument import load
    from odf.draw import Frame, Image
    import tempfile
    images = []
    with tempfile.NamedTemporaryFile(delete=True, suffix=".odp") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        for frame in doc.getElementsByType(Frame):
            for img in frame.getElementsByType(Image):
                href = img.getAttribute("href")
                if href and href.startswith("Pictures/"):
                    for f in doc.Pictures:
                        if f == href:
                            img_bytes = doc.Pictures[f]
                            images.append(base64.b64encode(img_bytes).decode("utf-8"))
    return images

def extract_images_ods(buffer: bytes) -> list[str]:
    from odf.opendocument import load
    from odf.draw import Frame, Image
    import tempfile
    images = []
    with tempfile.NamedTemporaryFile(delete=True, suffix=".ods") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        for frame in doc.getElementsByType(Frame):
            for img in frame.getElementsByType(Image):
                href = img.getAttribute("href")
                if href and href.startswith("Pictures/"):
                    for f in doc.Pictures:
                        if f == href:
                            img_bytes = doc.Pictures[f]
                            images.append(base64.b64encode(img_bytes).decode("utf-8"))
    return images

def extract_images_epub(buffer: bytes) -> list[str]:
    from ebooklib import epub
    import tempfile
    images = []
    with tempfile.NamedTemporaryFile(delete=True, suffix=".epub") as tmp:
        tmp.write(buffer)
        tmp.flush()
        book = epub.read_epub(tmp.name)
        for item in book.get_items():
            if item.get_type() == 3:  # ITEM_IMAGE
                images.append(base64.b64encode(item.get_content()).decode("utf-8"))
    return images

def extract_images_xlsx(buffer: bytes) -> list[str]:
    import tempfile
    import zipfile
    images = []
    with tempfile.NamedTemporaryFile(delete=True, suffix=".xlsx") as tmp:
        tmp.write(buffer)
        tmp.flush()
        # openpyxl only exposes images in worksheet._images, but not all images are there
        # fallback: extract from archive
        with zipfile.ZipFile(tmp.name) as zf:
            for name in zf.namelist():
                if name.startswith("xl/media/"):
                    img_bytes = zf.read(name)
                    images.append(base64.b64encode(img_bytes).decode("utf-8"))
    return images


def extract_text_docx_like(buffer: bytes) -> str:
    doc = Document(io.BytesIO(buffer))
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_mammoth(buffer: bytes) -> str:
    """
    Extract text from .dotx, .dotm, .docm using mammoth (supports more OOXML types).
    """
    import mammoth

    with io.BytesIO(buffer) as f:
        result = mammoth.extract_raw_text(f)
        return result.value.strip()


def extract_text_xlsx(buffer: bytes) -> str:
    """
    Extract text from .xlsx using openpyxl.
    """
    import openpyxl
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".xlsx") as tmp:
        tmp.write(buffer)
        tmp.flush()
        wb = openpyxl.load_workbook(tmp.name, data_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append("\t".join(row_text))
        return "\n".join(text)


def extract_text_xls(buffer: bytes) -> str:
    """
    Extract text from .xls using xlrd.
    """
    import xlrd
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".xls") as tmp:
        tmp.write(buffer)
        tmp.flush()
        wb = xlrd.open_workbook(tmp.name)
        text = []
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                row_text = [
                    str(cell.value) if cell.value is not None else "" for cell in row
                ]
                text.append("\t".join(row_text))
        return "\n".join(text)


def extract_text_odt(buffer: bytes) -> str:
    """
    Extract text from .odt using odfpy.
    """
    from odf.opendocument import load
    from odf.text import P
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".odt") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        paragraphs = doc.getElementsByType(P)
        return "\n".join([str(p) for p in paragraphs if p])


def extract_text_ods(buffer: bytes) -> str:
    """
    Extract text from .ods using odfpy.
    """
    from odf.opendocument import load
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".ods") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        tables = doc.getElementsByType(Table)
        text = []
        for table in tables:
            for row in table.getElementsByType(TableRow):
                row_text = []
                for cell in row.getElementsByType(TableCell):
                    cell_text = " ".join([str(p) for p in cell.getElementsByType(P)])
                    row_text.append(cell_text)
                text.append("\t".join(row_text))
        return "\n".join(text)


def extract_text_odp(buffer: bytes) -> str:
    """
    Extract text from .odp using odfpy.
    """
    from odf.opendocument import load
    from odf.draw import Page
    from odf.text import P
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".odp") as tmp:
        tmp.write(buffer)
        tmp.flush()
        doc = load(tmp.name)
        pages = doc.getElementsByType(Page)
        text = []
        for page in pages:
            for p in page.getElementsByType(P):
                text.append(str(p))
        return "\n".join(text)


def extract_text_pptx_like(buffer: bytes) -> str:
    """
    Extract text from .pptx, .ppsx, .pptm using python-pptx.
    """
    import io
    from pptx import Presentation

    with io.BytesIO(buffer) as f:
        prs = Presentation(f)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)


def extract_text_ppt(buffer: bytes) -> str:
    """
    Extract text from .ppt (legacy PowerPoint). Not natively supported by python-pptx.
    Raise error for now.
    """
    raise NotImplementedError(
        ".ppt extraction is not supported. Consider using unoconv or libreoffice for conversion."
    )


def extract_text_epub(buffer: bytes) -> str:
    """
    Extract text from .epub using ebooklib.
    """
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import tempfile

    with tempfile.NamedTemporaryFile(delete=True, suffix=".epub") as tmp:
        tmp.write(buffer)
        tmp.flush()
        book = epub.read_epub(tmp.name)
        text = []
        for item in book.get_items():
            # 9 is the ITEM_DOCUMENT constant in ebooklib
            if item.get_type() == 9:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text.append(soup.get_text(separator=" ", strip=True))
        return "\n".join(text)


def extract_text_ps(buffer: bytes) -> str:
    # Write to temp file, then use ps2ascii (Ghostscript) to extract text
    import tempfile

    try:
        with tempfile.NamedTemporaryFile(delete=True, suffix=".ps") as tmp:
            tmp.write(buffer)
            tmp.flush()
            result = subprocess.run(
                ["ps2ascii", tmp.name],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=True,
            )
            text = result.stdout.decode("utf-8", errors="replace")
        return text
    except FileNotFoundError:
        raise RuntimeError(
            "ps2ascii (Ghostscript) is not installed. Please install Ghostscript to extract text from .ps files."
        )
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from .ps file: {e}")
